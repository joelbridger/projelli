#!/usr/bin/env python3
"""
Generate the Keepance campaign fixture corpus.

Outputs into the parent directory (tests/fixtures/matter-corpus/).

Run from any location:
    python3 tests/fixtures/matter-corpus/generators/generate-fixtures.py

Regenerate a subset by naming makers (see MAKERS in main), e.g. only the
VG-2 scanned-PDF fixtures, leaving every other committed binary untouched:
    python3 .../generate-fixtures.py scanned-filing scanned-fax

Requires: python-docx, openpyxl, python-pptx (pip3 install python-docx openpyxl python-pptx)
Scanned-PDF fixtures additionally require Pillow (pip3 install pillow)
"""

import os
import sys
import struct
import zlib
import textwrap
from pathlib import Path

# ---------------------------------------------------------------------------
# Setup paths
# ---------------------------------------------------------------------------
SCRIPT_DIR = Path(__file__).parent
CORPUS_DIR = SCRIPT_DIR.parent  # tests/fixtures/matter-corpus/
MATTER_B_DIR = CORPUS_DIR / "matter-b-acme"
MATTER_B_DIR.mkdir(parents=True, exist_ok=True)

def out(name: str) -> Path:
    return CORPUS_DIR / name

def out_b(name: str) -> Path:
    return MATTER_B_DIR / name

# ---------------------------------------------------------------------------
# Import libraries
# ---------------------------------------------------------------------------
try:
    from docx import Document as DocxDocument
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("WARNING: python-docx not available", file=sys.stderr)

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False
    print("WARNING: openpyxl not available", file=sys.stderr)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PtPptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("WARNING: python-pptx not available", file=sys.stderr)

# VG-2 (Wave 2 Task 8) — scanned-PDF fixtures need Pillow (pip3 install pillow)
try:
    from PIL import Image, ImageDraw, ImageFont, ImageFilter
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    print("WARNING: Pillow not available (scanned PDF fixtures skipped)", file=sys.stderr)

# ---------------------------------------------------------------------------
# Helper: add tracked-change XML elements (w:ins / w:del) to a docx
# python-docx doesn't natively author tracked changes, so we build the OOXML
# directly and insert it into the document body XML.
# ---------------------------------------------------------------------------

def add_tracked_insert(doc, text: str, author: str, date: str, rev_id: int):
    """Append a w:ins tracked insertion paragraph to doc."""
    # Build the paragraph with a w:ins run
    p_elem = OxmlElement('w:p')
    ins = OxmlElement('w:ins')
    ins.set(qn('w:id'), str(rev_id))
    ins.set(qn('w:author'), author)
    ins.set(qn('w:date'), date)
    r = OxmlElement('w:r')
    t = OxmlElement('w:t')
    t.text = text
    r.append(t)
    ins.append(r)
    p_elem.append(ins)
    doc.element.body.append(p_elem)

def add_tracked_delete(doc, text: str, author: str, date: str, rev_id: int):
    """Append a w:del tracked deletion paragraph to doc."""
    p_elem = OxmlElement('w:p')
    del_elem = OxmlElement('w:del')
    del_elem.set(qn('w:id'), str(rev_id))
    del_elem.set(qn('w:author'), author)
    del_elem.set(qn('w:date'), date)
    dr = OxmlElement('w:r')
    dt = OxmlElement('w:delText')
    dt.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')
    dt.text = text
    dr.append(dt)
    del_elem.append(dr)
    p_elem.append(del_elem)
    doc.element.body.append(p_elem)

def add_comment(doc, paragraph_index: int, comment_text: str, comment_author: str,
                comment_date: str, comment_id: int):
    """Add a comment anchor to a paragraph (body-side only).

    Full comments.xml injection is done at the OOXML level in post_process_add_comments().
    This function adds the comment range markers to the paragraph.
    """
    body = doc.element.body
    paragraphs = [ch for ch in body if ch.tag.endswith('}p')]
    if paragraph_index < len(paragraphs):
        target_p = paragraphs[paragraph_index]
        # Prepend range start
        rs = OxmlElement('w:commentRangeStart')
        rs.set(qn('w:id'), str(comment_id))
        target_p.insert(0, rs)
        # Append range end + reference
        re_ = OxmlElement('w:commentRangeEnd')
        re_.set(qn('w:id'), str(comment_id))
        target_p.append(re_)
        ref_r = OxmlElement('w:r')
        rpr = OxmlElement('w:rPr')
        rpr_style = OxmlElement('w:rStyle')
        rpr_style.set(qn('w:val'), 'CommentReference')
        rpr.append(rpr_style)
        ref_r.append(rpr)
        cref = OxmlElement('w:commentReference')
        cref.set(qn('w:id'), str(comment_id))
        ref_r.append(cref)
        target_p.append(ref_r)


# ---------------------------------------------------------------------------
# 1. engagement-letter-tracked.docx
#    2-page engagement letter with 4 tracked changes (2 ins, 2 del) by 2
#    authors, plus 2 comments.
# ---------------------------------------------------------------------------
def make_engagement_letter_tracked():
    if not HAS_DOCX:
        print("SKIP: engagement-letter-tracked.docx (no python-docx)")
        return
    path = out("engagement-letter-tracked.docx")
    doc = DocxDocument()
    doc.add_heading("ENGAGEMENT LETTER", 0)
    doc.add_paragraph("Marchetti & Associates LLP")
    doc.add_paragraph("1200 Commerce Drive, Suite 400")
    doc.add_paragraph("New York, NY 10001")
    doc.add_paragraph("")
    doc.add_paragraph("June 10, 2026")
    doc.add_paragraph("")
    doc.add_paragraph("Re: Retention — Johnson v. Nexus Dynamics Corp.")
    doc.add_paragraph("")
    p1 = doc.add_paragraph(
        "Dear Mr. Johnson, this letter confirms the terms of our engagement. "
        "We will represent you in connection with your wrongful termination claim "
        "against Nexus Dynamics Corp. Our hourly rate is $450 per hour."
    )
    doc.add_paragraph("")
    p2 = doc.add_paragraph(
        "Scope of Representation: We will handle all aspects of pre-litigation "
        "negotiation, mediation, and if necessary, litigation through trial. "
        "We will not handle any appeals without a separate engagement agreement."
    )
    doc.add_paragraph("")
    doc.add_paragraph(
        "Retainer: You agree to pay an initial retainer of $10,000. "
        "Additional costs will be billed monthly."
    )
    doc.add_paragraph("")
    doc.add_paragraph(
        "Confidentiality: All communications between attorney and client are "
        "protected by attorney-client privilege. You must not share our work "
        "product with third parties without prior written consent."
    )
    doc.add_paragraph("")
    # Page 2 content
    doc.add_page_break()
    doc.add_heading("TERMS AND CONDITIONS", 1)
    doc.add_paragraph(
        "Governing Law: This agreement shall be governed by the laws of the "
        "State of New York."
    )
    doc.add_paragraph(
        "Dispute Resolution: Any fee disputes shall first be submitted to "
        "non-binding mediation before the New York State Bar Association."
    )
    doc.add_paragraph(
        "Termination: Either party may terminate this engagement upon ten (10) "
        "days written notice. Client remains responsible for all fees incurred "
        "through the termination date."
    )
    doc.add_paragraph("")
    doc.add_paragraph(
        "Please sign below to acknowledge your agreement to these terms."
    )
    doc.add_paragraph("")
    doc.add_paragraph("Sincerely,")
    doc.add_paragraph("")
    doc.add_paragraph("Diane Marchetti, Esq.")
    doc.add_paragraph("Marchetti & Associates LLP")
    doc.add_paragraph("")
    doc.add_paragraph("Accepted by: ________________________  Date: __________")

    # Add 4 tracked changes (2 authors: "Diane Marchetti" and "James Thornton")
    # Tracked insertion 1 — Marchetti adds a clause
    add_tracked_insert(
        doc,
        "[INSERTED BY MARCHETTI] Additionally, we will provide monthly status updates.",
        "Diane Marchetti",
        "2026-06-05T10:00:00Z",
        rev_id=1,
    )
    # Tracked deletion 1 — Marchetti deletes placeholder
    add_tracked_delete(
        doc,
        "[TO BE REMOVED: placeholder text from template]",
        "Diane Marchetti",
        "2026-06-05T10:05:00Z",
        rev_id=2,
    )
    # Tracked insertion 2 — Thornton (opposing review) adds comment about rate
    add_tracked_insert(
        doc,
        "[INSERTED BY THORNTON] Hourly rate subject to annual review.",
        "James Thornton",
        "2026-06-06T09:00:00Z",
        rev_id=3,
    )
    # Tracked deletion 2 — Thornton removes duplicated sentence
    add_tracked_delete(
        doc,
        "[DELETED BY THORNTON: duplicated clause — see section 4]",
        "James Thornton",
        "2026-06-06T09:10:00Z",
        rev_id=4,
    )

    # Add comment anchors in the body (range markers)
    try:
        add_comment(
            doc,
            paragraph_index=0,
            comment_text="Confirm client identity before sending. Do not share until retainer received.",
            comment_author="Diane Marchetti",
            comment_date="2026-06-05T10:30:00Z",
            comment_id=1,
        )
        add_comment(
            doc,
            paragraph_index=8,
            comment_text="Verify retainer amount against current firm policy — was $10k as of Jan 2026.",
            comment_author="James Thornton",
            comment_date="2026-06-06T09:15:00Z",
            comment_id=2,
        )
    except Exception as e:
        print(f"  Note: comment anchor injection error (non-fatal): {e}")

    # Save to a temp path first, then inject comments.xml via zipfile surgery
    import tempfile, zipfile, shutil
    tmp = path.with_suffix(".tmp.docx")
    doc.save(str(tmp))

    # Now inject comments.xml into the saved zip
    comments_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:comments xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        '<w:comment w:id="1" w:author="Diane Marchetti" w:date="2026-06-05T10:30:00Z" w:initials="DM">'
        '<w:p><w:r><w:t>Confirm client identity before sending. Do not share until retainer received.</w:t></w:r></w:p>'
        '</w:comment>'
        '<w:comment w:id="2" w:author="James Thornton" w:date="2026-06-06T09:15:00Z" w:initials="JT">'
        '<w:p><w:r><w:t>Verify retainer amount against current firm policy — was $10k as of Jan 2026.</w:t></w:r></w:p>'
        '</w:comment>'
        '</w:comments>'
    )

    COMMENTS_REL_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/comments'
    COMMENTS_CT = 'application/vnd.openxmlformats-officedocument.wordprocessingml.comments+xml'

    try:
        with zipfile.ZipFile(str(tmp), 'r') as zin, zipfile.ZipFile(str(path), 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename == '[Content_Types].xml':
                    # Add comments content type
                    ct_str = data.decode('utf-8')
                    if 'comments.xml' not in ct_str:
                        ct_str = ct_str.replace(
                            '</Types>',
                            f'<Override PartName="/word/comments.xml" ContentType="{COMMENTS_CT}"/></Types>'
                        )
                    zout.writestr(item, ct_str.encode('utf-8'))
                elif item.filename == 'word/_rels/document.xml.rels':
                    # Add comments relationship
                    rels_str = data.decode('utf-8')
                    if 'comments' not in rels_str:
                        rels_str = rels_str.replace(
                            '</Relationships>',
                            f'<Relationship Id="rId99" Type="{COMMENTS_REL_TYPE}" Target="comments.xml"/></Relationships>'
                        )
                    zout.writestr(item, rels_str.encode('utf-8'))
                else:
                    zout.writestr(item, data)
            # Add the comments.xml part
            zout.writestr('word/comments.xml', comments_xml.encode('utf-8'))
        tmp.unlink()
        print(f"CREATED: {path.name} (with 4 tracked changes + 2 comments)")
    except Exception as e:
        # Fall back to plain save without comments injection
        shutil.move(str(tmp), str(path))
        print(f"CREATED: {path.name} (comments injection failed: {e}; saved without comments.xml)")


# ---------------------------------------------------------------------------
# 2. deposition-transcript-johnson.txt
#    120+ line Q/A deposition transcript. 3 planted contradictions vs
#    incident-summary-johnson.md (labelled CONTRADICTION-1/2/3 in comments).
# ---------------------------------------------------------------------------
TRANSCRIPT = """\
IN THE UNITED STATES DISTRICT COURT
FOR THE SOUTHERN DISTRICT OF NEW YORK

-------------------------------------------------------
IN RE: JOHNSON v. NEXUS DYNAMICS CORP.
Case No. 26-CV-0412 (S.D.N.Y.)
-------------------------------------------------------

DEPOSITION OF MARCUS JOHNSON
Taken on May 14, 2026 at 9:00 a.m.
Offices of Marchetti & Associates LLP
1200 Commerce Drive, Suite 400, New York, NY 10001

-------------------------------------------------------
APPEARANCES:
  FOR THE PLAINTIFF:   Diane Marchetti, Esq.
  FOR THE DEFENDANT:   Kevin Hartley, Esq. (Nexus Dynamics Corp.)
  COURT REPORTER:      Sandra Price, RPR
-------------------------------------------------------

EXAMINATION BY MS. MARCHETTI:

PAGE 1

Q. Please state your full name for the record.
A. Marcus Raymond Johnson.

Q. And your current address?
A. 47 Riverside Drive, Apartment 12B, New York, New York 10024.

Q. Mr. Johnson, how long were you employed by Nexus Dynamics Corp.?
A. Six years and three months. I started in February 2020.

Q. What was your title at the time of your termination?
A. Senior Product Manager, Cloud Infrastructure Division.

Q. And when were you terminated?
A. November 12, 2025.

Q. Did you receive any written notice of the termination?
A. No. They called me into HR and told me verbally.

Q. Who was present in that meeting?
A. Sandra Liu from HR, and my direct manager, Tom Weston.

Q. Did Mr. Weston say anything about the reason for your termination?
A. He said it was a "restructuring," but he didn't give specifics.

Q. Had you received any performance reviews in the six months prior?
A. Yes, my last review was in August 2025. I received a "Meets Expectations" rating.

Q. Was there any mention of performance concerns?
A. No. My August review had no written performance concerns.

PAGE 2

Q. Let me direct your attention to September 2025. Were you aware of
   an internal investigation by Nexus Dynamics at that time?
A. No, I was not aware of any investigation involving me.

Q. Do you recall receiving an email from Compliance on September 8, 2025?
A. I receive many emails. I don't specifically recall one from Compliance
   in September.

Q. I am going to show you what has been marked as Plaintiff's Exhibit 3.
   Is this an email you received?
A. Yes, that appears to be an email to me from the Compliance department.

Q. And what is the date on that email?
A. September 8, 2025.

Q. And what does it say?
A. It says there is a review of expense reports from Q2 2025 and asks
   me to preserve all related documents.

Q. Did you preserve those documents?
A. I believe I did. I forwarded them to my personal email for safekeeping.
   [CONTRADICTION-1: Transcript says he forwarded to personal email for
   safekeeping. The incident summary states he did NOT forward any documents
   outside company systems and all materials remained on company servers only.]

Q. When exactly did you forward those documents?
A. I think it was September 9 or 10. Within a day or two of receiving
   the Compliance email.

Q. Did anyone from Nexus Dynamics ever tell you that forwarding documents
   to personal email was a policy violation?
A. Not explicitly. I was not aware of any policy against it.

PAGE 3

Q. Let me turn to the incident of October 3, 2025. Were you in the
   office that day?
A. Yes.

Q. Do you recall a meeting with Mr. Weston at approximately 2:00 p.m.?
A. Yes, I recall a meeting.

Q. What was the meeting about?
A. Tom told me the Q2 expense review had found some discrepancies and
   that I would need to submit a written explanation.

Q. How did you respond?
A. I told him I was happy to cooperate and that I had already saved
   all the relevant documentation.

Q. Did Mr. Weston tell you a deadline for submitting the explanation?
A. He said I had until October 17, 2025 to submit my written response.
   [CONTRADICTION-2: Transcript records the deadline as October 17, 2025.
   The incident summary states the deadline given to Johnson was October 10, 2025.]

Q. Did you meet that deadline?
A. I submitted my response on October 15, 2025, two days before the deadline.

Q. Let's talk about your relationship with Mr. Weston. Would you describe
   it as positive?
A. Mostly yes. We had occasional disagreements about project priorities,
   but nothing serious.

PAGE 4

Q. In the eighteen months before your termination, did you ever raise any
   concerns about Nexus Dynamics' business practices?
A. Yes. In March 2025 I raised concerns about billing irregularities I
   had noticed in vendor invoices.

Q. Who did you raise those concerns with?
A. I sent an email to the ethics hotline and also spoke directly with
   Sandra Liu in HR.

Q. And what response did you receive?
A. Sandra said she would look into it. I never received a written follow-up.

Q. Did you keep a copy of that email?
A. Yes, I have that email.

Q. After you raised those concerns, did you notice any change in how
   Mr. Weston treated you?
A. Yes. Starting in April 2025, Tom began excluding me from key meetings
   that I had previously attended. My responsibilities were also reduced.

Q. How many meetings do you estimate you were excluded from between
   April and October 2025?
A. I would say at least fifteen to twenty significant meetings.

PAGE 5

Q. Mr. Johnson, I want to ask you about your compensation. What was your
   base salary at the time of termination?
A. $185,000 per year.

Q. Did you receive any bonuses in 2024?
A. Yes. I received my standard annual bonus of $22,000 in March 2024.

Q. And for 2025?
A. I did not receive a bonus for 2025. I was told bonuses were suspended
   company-wide due to budget constraints.

Q. Were you aware of other employees receiving bonuses in 2025?
A. Yes. I learned from colleagues that two other product managers received
   retention bonuses in September 2025.

Q. Do you know the names of those colleagues?
A. I'd prefer not to put their names on record, but yes, I know who they are.

Q. At the time of your termination, did anyone at Nexus Dynamics explain
   the severance package being offered?
A. Sandra Liu gave me a document describing a four-week severance.
   [CONTRADICTION-3: Transcript says he was offered four weeks of severance.
   The incident summary states the offer was for eight weeks of severance per
   company policy for employees with more than five years of service.]

Q. Did you sign any severance agreement?
A. No. I declined to sign because I wanted to consult with an attorney first.

Q. Is there anything else you believe is important for the record at
   this time?
A. Only that I believe my termination was directly connected to my
   raising concerns about vendor billing in March 2025.

MS. MARCHETTI: Thank you, Mr. Johnson. No further questions at this time.

EXAMINATION BY MR. HARTLEY:

Q. Mr. Johnson, you mentioned forwarding documents to your personal email.
   Were you aware that Nexus Dynamics' IT policy requires all work documents
   to remain on company-managed systems?
A. I may have seen something about that policy, but I didn't consider it
   relevant to preserving documents for a compliance review.

Q. And regarding the October meeting with Mr. Weston — you said the
   deadline was October 17. Do you have any written confirmation of that date?
A. Tom told me verbally. I didn't receive it in writing.

MR. HARTLEY: No further questions.

MS. MARCHETTI: We'll reserve the right to re-examine.

-------------------------------------------------------
CERTIFICATE OF REPORTER
I, Sandra Price, RPR, Certified Court Reporter, do hereby certify that
the foregoing is a true and accurate transcript of the deposition of
Marcus Johnson taken on May 14, 2026.

Sandra Price, RPR
-------------------------------------------------------
"""

def make_deposition_transcript():
    path = out("deposition-transcript-johnson.txt")
    path.write_text(TRANSCRIPT, encoding="utf-8")
    print(f"CREATED: {path.name} ({len(TRANSCRIPT.splitlines())} lines)")


# ---------------------------------------------------------------------------
# 2b. deposition-transcript-weston-certified.txt (VG-3c, Wave 2 Task 9)
#    The CERTIFIED line-numbered deposition format: every content line carries
#    a right-aligned two-digit line number 1-25, pages separated by a centered
#    "Page N" header. detect_transcript (src-tauri rag/transcript.rs) must
#    return TRUE for this file and FALSE for deposition-transcript-johnson.txt
#    (which is deliberately Q./A. + "PAGE N" prose, NOT line-numbered).
#
#    Content: document-retention and litigation-hold PROCEDURE testimony by
#    Thomas Weston — topically NEUTRAL: it must not corroborate or contradict
#    any of the three planted contradiction pairs (personal-email forwarding,
#    the written-response deadline, the severance weeks), so the finder
#    rubric dynamics stay clean.
# ---------------------------------------------------------------------------

# The unique quotable sentence planted at page 2 lines 14-16 (spoken words,
# reconstructed across the transcript's line wraps). The leg-1 retrieval test
# (src-tauri/tests/rag_deposition_contradictions.rs
# certified_transcript_chunks_carry_page_line_locators) cites these constants
# verbatim — change them ONLY together with that test.
WESTON_HOLD_SENTENCE = (
    "The litigation hold notice went out to the cloud infrastructure team "
    "on September 12, 2025."
)
WESTON_HOLD_PAGE = 2
WESTON_HOLD_LINES = (14, 16)

# 25 content lines per page, exactly. Line numbers are applied by the maker.
WESTON_PAGES = [
    [  # Page 1 — caption + examination opens
        "UNITED STATES DISTRICT COURT",
        "SOUTHERN DISTRICT OF NEW YORK",
        "------------------------------------------x",
        "IN RE: JOHNSON v. NEXUS DYNAMICS CORP.",
        "Case No. 26-CV-0412 (S.D.N.Y.)",
        "------------------------------------------x",
        "DEPOSITION OF THOMAS WESTON",
        "Taken on June 3, 2026 at 9:30 a.m.",
        "Offices of Marchetti & Associates LLP",
        "1200 Commerce Drive, Suite 400, New York",
        "APPEARANCES:",
        "For the Plaintiff: Diane Marchetti, Esq.",
        "For the Defendant: Kevin Hartley, Esq.",
        "Court Reporter: Sandra Price, RPR",
        "EXAMINATION BY MS. MARCHETTI:",
        "Q.   Please state your name for the record.",
        "A.   Thomas Edward Weston.",
        "Q.   What is your position at Nexus Dynamics?",
        "A.   Senior Director of the Cloud",
        "Infrastructure Division.",
        "Q.   How long have you held that role?",
        "A.   Since March of 2019.",
        "Q.   Are you familiar with the company's",
        "document retention policy?",
        "A.   Yes, I am. I helped draft the current",
    ],
    [  # Page 2 — retention schedule + the planted hold-notice lines 14-16
        "version of that policy in 2023.",
        "Q.   Describe the retention schedule for",
        "engineering records, please.",
        "A.   Project documentation is retained for",
        "seven years. Routine operational logs are",
        "kept for eighteen months and then purged",
        "on a rolling basis by the archive system.",
        "Q.   Who administers litigation holds at",
        "Nexus Dynamics?",
        "A.   The legal operations group issues the",
        "hold notices. My division carries out the",
        "preservation steps for cloud systems.",
        "Q.   When was a hold issued for this case?",
        "A.   The litigation hold notice went out to",
        "the cloud infrastructure team on September",
        "12, 2025.",
        "Q.   What does the team do when a notice",
        "like that arrives?",
        "A.   Automated purges are suspended for the",
        "named custodians and the backup snapshots",
        "are moved to the legal archive tier.",
        "Q.   Is that suspension logged anywhere?",
        "A.   Yes. The archive system writes an",
        "entry for every suspension and every",
        "release. The log is append-only.",
    ],
    [  # Page 3 — backups + training (neutral filler)
        "Q.   Turning to backups. How are the cloud",
        "system backups organized?",
        "A.   Nightly incremental snapshots and a",
        "full weekly snapshot, retained on a",
        "ninety-day rotation unless a hold applies.",
        "Q.   Where are the snapshots stored?",
        "A.   In a separate archive account with",
        "restricted access. Two administrators in",
        "my division hold the credentials.",
        "Q.   Does the retention policy cover email?",
        "A.   Email is managed by corporate IT, not",
        "my division. My understanding is that the",
        "same hold process applies to mailboxes.",
        "Q.   How are employees trained on these",
        "procedures?",
        "A.   There is an annual compliance module.",
        "Completion is tracked by Human Resources",
        "and reported to the audit committee.",
        "Q.   Did your division complete the 2025",
        "module?",
        "A.   Yes. Completion for my division was",
        "one hundred percent as of October 2025.",
        "Q.   Thank you. Let's take a short break.",
        "(Recess taken from 10:42 a.m. to",
        "10:55 a.m.)",
    ],
    [  # Page 4 — cross + certificate
        "EXAMINATION BY MR. HARTLEY:",
        "Q.   Mr. Weston, just a few questions. Is",
        "the retention policy you described written",
        "down anywhere?",
        "A.   Yes, it is published on the internal",
        "policy portal as document RET-104.",
        "Q.   And the archive log you mentioned, who",
        "can read it?",
        "A.   Legal operations and the two archive",
        "administrators. Read access is logged as",
        "well.",
        "Q.   Nothing further. Thank you.",
        "MS. MARCHETTI: No redirect. We are",
        "adjourned.",
        "(Deposition concluded at 11:20 a.m.)",
        "CERTIFICATE OF REPORTER",
        "I, Sandra Price, RPR, Certified Court",
        "Reporter, do hereby certify that the",
        "foregoing is a true and accurate",
        "transcript of the deposition of Thomas",
        "Weston, taken on June 3, 2026.",
        "Sandra Price, RPR",
        "Notary Public, State of New York",
        "Commission No. 01PR6054321",
        "My commission expires March 14, 2028.",
    ],
]


def make_weston_certified_transcript():
    path = out("deposition-transcript-weston-certified.txt")
    out_lines = []
    for page_no, lines in enumerate(WESTON_PAGES, start=1):
        assert len(lines) == 25, f"page {page_no} must have exactly 25 lines"
        header = f"Page {page_no}"
        out_lines.append(" " * max(0, (58 - len(header)) // 2) + header)
        out_lines.append("")
        for line_no, content in enumerate(lines, start=1):
            out_lines.append(f"{line_no:>2}   {content}")
        out_lines.append("")
    # Self-check: the planted sentence sits exactly where the constants say
    # (whitespace-normalized across the line wraps, the same canon the
    # citation verifier uses).
    lo, hi = WESTON_HOLD_LINES
    planted = WESTON_PAGES[WESTON_HOLD_PAGE - 1][lo - 1 : hi]
    joined = " ".join(" ".join(planted).split())
    assert WESTON_HOLD_SENTENCE in joined, f"planted sentence drifted: {joined!r}"
    text = "\n".join(out_lines) + "\n"
    path.write_text(text, encoding="utf-8")
    print(
        f"CREATED: {path.name} ({len(WESTON_PAGES)} certified pages, "
        f"hold-notice sentence at page {WESTON_HOLD_PAGE} "
        f"lines {WESTON_HOLD_LINES[0]}-{WESTON_HOLD_LINES[1]})"
    )


# ---------------------------------------------------------------------------
# 3. incident-summary-johnson.md
#    Matter summary that contradicts the transcript in 3 places (labelled).
# ---------------------------------------------------------------------------
INCIDENT_SUMMARY = """\
# Matter Summary: Johnson v. Nexus Dynamics Corp.

**Case No.:** 26-CV-0412 (S.D.N.Y.)
**Client:** Marcus Raymond Johnson
**Opposing Party:** Nexus Dynamics Corp.
**Lead Attorney:** Diane Marchetti, Esq.
**Date Prepared:** June 1, 2026

---

## Background

Marcus Johnson was employed by Nexus Dynamics Corp. as a Senior Product Manager
from February 2020 until his termination on November 12, 2025 — a tenure of
approximately five years and nine months.

Johnson's termination followed an internal compliance review of Q2 2025 expense
reports. The company characterizes the separation as a "restructuring." Johnson
contends it was retaliatory, stemming from his March 2025 report of suspected
vendor billing irregularities to the company ethics hotline.

---

## Key Facts (Established from Company Records)

### Document Preservation

Following receipt of the Compliance Department's document preservation notice
dated September 8, 2025, Johnson confirmed to HR that **all relevant documents
remained on company servers only** and that he did not transfer or copy any
materials to personal devices or email accounts.

> [CONTRADICTION-1 vs. deposition transcript: The transcript records Johnson
> saying he forwarded documents to his personal email for safekeeping on
> approximately September 9-10. This summary states no documents left company
> systems. These accounts directly conflict on a material fact.]

### October 3 Meeting — Written Response Deadline

During the October 3, 2025 meeting with Tom Weston (manager) and Sandra Liu (HR),
Johnson was given a deadline of **October 10, 2025** to submit his written
explanation of the Q2 expense discrepancies.

> [CONTRADICTION-2 vs. deposition transcript: The transcript records the
> deadline as October 17, 2025. The incident summary records it as October 10.
> A seven-day discrepancy on a key compliance deadline is directly material.]

### Severance Offer

Upon termination, Nexus Dynamics presented Johnson with a severance agreement
offering **eight (8) weeks of base salary continuation** in accordance with
company policy for employees with five or more years of service.

> [CONTRADICTION-3 vs. deposition transcript: The transcript records Johnson
> stating he was offered four weeks of severance. The incident summary records
> the offer as eight weeks. This discrepancy affects both damages calculation
> and the question of whether Johnson understood the offer presented.]

---

## Timeline of Key Events

| Date | Event |
|------|-------|
| Feb 2020 | Johnson hired as Product Manager |
| Mar 2025 | Johnson reports billing irregularities to ethics hotline and HR |
| Apr 2025 | Johnson excluded from meetings; reduced responsibilities (per Johnson) |
| Aug 2025 | Last performance review: "Meets Expectations," no written concerns |
| Sep 8, 2025 | Compliance sends document preservation notice |
| Oct 3, 2025 | Meeting with Weston and Liu regarding Q2 expense discrepancies |
| Oct 10, 2025 | Deadline for Johnson's written response (per this summary) |
| Oct 15, 2025 | Johnson submits written response |
| Nov 12, 2025 | Johnson terminated |

---

## Litigation Posture

**Theory of recovery:** Retaliatory termination under New York Labor Law § 740
(whistleblower protection) and Title VII (if discriminatory motive can be shown).

**Key vulnerabilities:** Document forwarding to personal email (IT policy
compliance); possible gap between his subjective belief about retaliation and
documentary proof of causal nexus.

**Recommended next steps:**
1. Resolve the document-forwarding contradiction before any motion practice.
2. Obtain written confirmation of the October compliance deadline from Weston
   or Liu.
3. Subpoena Nexus Dynamics for September–November 2025 manager communications.
4. Identify and interview the two Product Managers who received retention
   bonuses in September 2025.
"""

def make_incident_summary():
    path = out("incident-summary-johnson.md")
    path.write_text(INCIDENT_SUMMARY, encoding="utf-8")
    print(f"CREATED: {path.name}")


# ---------------------------------------------------------------------------
# 4. contract-services-agreement.docx — simple 5-clause contract
# ---------------------------------------------------------------------------
def make_services_agreement():
    if not HAS_DOCX:
        print("SKIP: contract-services-agreement.docx (no python-docx)")
        return
    path = out("contract-services-agreement.docx")
    doc = DocxDocument()
    doc.add_heading("PROFESSIONAL SERVICES AGREEMENT", 0)
    doc.add_paragraph("This Professional Services Agreement (\"Agreement\") is entered into as of June 10, 2026, by and between Marchetti & Associates LLP (\"Firm\") and Nexus Consulting Partners LLC (\"Client\").")
    doc.add_paragraph("")
    doc.add_heading("1. Scope of Services", 1)
    doc.add_paragraph("The Firm agrees to provide legal research, document review, and litigation support services as requested by Client in writing. All services shall be performed in accordance with applicable professional standards.")
    doc.add_heading("2. Fees and Payment", 1)
    doc.add_paragraph("Client agrees to pay the Firm at a blended hourly rate of $375 per hour. Invoices shall be rendered monthly and are due within thirty (30) days of receipt. Late payments accrue interest at 1.5% per month.")
    doc.add_heading("3. Confidentiality", 1)
    doc.add_paragraph("Each party agrees to maintain in strict confidence all Confidential Information received from the other party and to use such information only in furtherance of this Agreement. This obligation survives termination.")
    doc.add_heading("4. Term and Termination", 1)
    doc.add_paragraph("This Agreement commences on the date first written above and continues for twelve (12) months, unless earlier terminated. Either party may terminate upon thirty (30) days written notice. Upon termination, Client shall pay all fees incurred through the effective termination date.")
    doc.add_heading("5. Governing Law", 1)
    doc.add_paragraph("This Agreement shall be governed by and construed in accordance with the laws of the State of New York, without regard to its conflict of laws principles. Any disputes shall be resolved exclusively in the state or federal courts located in New York County.")
    doc.add_paragraph("")
    doc.add_paragraph("IN WITNESS WHEREOF, the parties have executed this Agreement as of the date first written above.")
    doc.add_paragraph("")
    doc.add_paragraph("MARCHETTI & ASSOCIATES LLP")
    doc.add_paragraph("By: ________________________")
    doc.add_paragraph("Name: Diane Marchetti, Esq.")
    doc.add_paragraph("")
    doc.add_paragraph("NEXUS CONSULTING PARTNERS LLC")
    doc.add_paragraph("By: ________________________")
    doc.add_paragraph("Name: ________________________")
    doc.save(str(path))
    print(f"CREATED: {path.name}")


# ---------------------------------------------------------------------------
# 5. scanned-exhibit.pdf — minimal valid 1-page PDF (hand-built bytes)
# ---------------------------------------------------------------------------
def make_scanned_exhibit():
    path = out("scanned-exhibit.pdf")
    # Minimal valid PDF with one page and a text object
    pdf_content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 120 >>
stream
BT
/F1 12 Tf
72 720 Td
(EXHIBIT A - SCANNED DOCUMENT) Tj
0 -20 Td
(Johnson v. Nexus Dynamics Corp. - Case 26-CV-0412) Tj
0 -20 Td
(This document is a placeholder exhibit for testing.) Tj
ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000274 00000 n
0000000446 00000 n
trailer
<< /Size 6 /Root 1 0 R >>
startxref
525
%%EOF
"""
    path.write_bytes(pdf_content)
    print(f"CREATED: {path.name}")


# ---------------------------------------------------------------------------
# 6. damages-model.xlsx
# ---------------------------------------------------------------------------
def make_damages_model():
    if not HAS_XLSX:
        print("SKIP: damages-model.xlsx (no openpyxl)")
        return
    path = out("damages-model.xlsx")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Damages"
    # Headers
    ws.append(["Category", "Amount ($)", "Notes"])
    # Data rows
    ws.append(["Lost wages (2025 Q4)", 46250.00, "Base $185k/4 = $46,250 for Nov 12 - Dec 31"])
    ws.append(["Lost wages (2026 projected)", 185000.00, "Full year base salary"])
    ws.append(["Lost bonus (2025)", 22000.00, "Standard annual bonus not paid"])
    ws.append(["Retention bonus (comparable employees)", 15000.00, "Estimate — two PM peers received"])
    ws.append(["Benefits lost (health insurance)", 12000.00, "Estimated annual employer contribution"])
    ws.append(["Emotional distress", 75000.00, "Subject to expert testimony"])
    ws.append(["Punitive damages (if awarded)", 500000.00, "Discretionary — max under NY law"])
    ws.append(["", "", ""])
    ws.append(["TOTAL (base)", "=SUM(B2:B7)", "Excludes punitive"])
    ws.append(["TOTAL (with punitive)", "=SUM(B2:B8)", "Full claim"])
    wb.save(str(path))
    print(f"CREATED: {path.name}")


# ---------------------------------------------------------------------------
# 7. exhibit-deck.pptx
# ---------------------------------------------------------------------------
def make_exhibit_deck():
    if not HAS_PPTX:
        print("SKIP: exhibit-deck.pptx (no python-pptx)")
        return
    path = out("exhibit-deck.pptx")
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    # Slide 1: Title
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Johnson v. Nexus Dynamics Corp."
    slide1.placeholders[1].text = "Evidence Exhibit Deck\nCase 26-CV-0412 | June 2026"
    # Slide 2: Timeline
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Key Events Timeline"
    slide2.placeholders[1].text = (
        "March 2025: Whistleblower report filed\n"
        "September 2025: Document preservation notice\n"
        "October 2025: Compliance meeting & deadline\n"
        "November 12, 2025: Termination"
    )
    prs.save(str(path))
    print(f"CREATED: {path.name}")


# ---------------------------------------------------------------------------
# 8. huge-notes.md — ~2MB of legal-flavored markdown
# ---------------------------------------------------------------------------
def make_huge_notes():
    path = out("huge-notes.md")
    block = textwrap.dedent("""\
    ## Litigation Note — Block {i}

    **Date:** June {day}, 2026
    **Prepared by:** Diane Marchetti, Esq.
    **Matter:** Johnson v. Nexus Dynamics Corp.

    ### Research Notes

    This section documents preliminary research into New York Labor Law § 740,
    the state whistleblower protection statute. Under § 740, an employer may not
    take retaliatory personnel action against an employee who discloses, threatens
    to disclose, objects to, or refuses to participate in an activity, policy, or
    practice of the employer that is in violation of law, rule, or regulation.

    The threshold question is whether Johnson's March 2025 report to the ethics
    hotline constitutes a "disclosure" within the meaning of § 740(2)(a). Case law
    in the Second Circuit suggests internal reporting to a company ethics hotline
    can satisfy the disclosure requirement where the employer subsequently takes
    adverse action. See *Pipia v. Nassau Cnty.*, 34 N.Y.3d 350 (2019).

    However, a causal nexus must be shown between the protected activity and the
    adverse employment action. Temporal proximity — here, approximately eight
    months between the March 2025 report and the November 2025 termination — may
    be insufficient on its own. The court in *Kotcher v. Rosa and Sullivan
    Appliance Ctr.*, 957 F.2d 59, 65 (2d Cir. 1992), held that a time gap of
    eight to ten months requires additional circumstantial evidence.

    Additional evidence to pursue:
    - Documentary evidence of Weston's changed conduct after April 2025
    - Email communications between Weston and Liu re: Johnson during Q2-Q4 2025
    - Comparator data: how were other employees' expense discrepancies handled?
    - Whether the "restructuring" justification was applied to any other PM

    ### Document Review Notes

    Plaintiff's Exhibit 3 (Compliance email of September 8, 2025) establishes
    that Nexus Dynamics had an ongoing investigation as of that date. The
    document preservation directive creates an inference that management knew
    of Johnson's potential evidentiary significance well before the termination.

    Defendant will likely argue that the termination decision was independent of
    the compliance review. We need to identify the decision-makers and their
    communications timeline more precisely.

    ### Privilege Log Notes

    Communications between Johnson and this firm are protected by attorney-client
    privilege. Work product prepared by or at the direction of counsel is protected
    under the work product doctrine. Marchetti & Associates has established a
    privilege log to track all potentially protected communications.

    """)

    chunks = []
    target_bytes = 2 * 1024 * 1024  # 2MB
    i = 0
    day = 1
    while sum(len(c.encode("utf-8")) for c in chunks) < target_bytes:
        chunks.append(block.format(i=i+1, day=(day % 28) + 1))
        i += 1
        day += 1

    content = "# Litigation Notes — Johnson v. Nexus Dynamics Corp.\n\n" + "".join(chunks)
    path.write_text(content, encoding="utf-8")
    size_mb = path.stat().st_size / (1024 * 1024)
    print(f"CREATED: {path.name} ({size_mb:.2f} MB)")


# ---------------------------------------------------------------------------
# 9. empty.docx — zero-byte file (deliberately invalid)
# ---------------------------------------------------------------------------
def make_empty_docx():
    path = out("empty.docx")
    path.write_bytes(b"")
    print(f"CREATED: {path.name} (0 bytes — intentionally invalid)")


# ---------------------------------------------------------------------------
# 10. Müller — Schäfer engagement (draft 2).docx — unicode/spaces/parens
# ---------------------------------------------------------------------------
def make_unicode_docx():
    if not HAS_DOCX:
        print("SKIP: Müller — Schäfer engagement (draft 2).docx (no python-docx)")
        return
    path = out("Müller — Schäfer engagement (draft 2).docx")
    doc = DocxDocument()
    doc.add_heading("ENGAGEMENT LETTER", 0)
    doc.add_paragraph("Müller & Schäfer GmbH")
    doc.add_paragraph("Hauptstraße 42")
    doc.add_paragraph("80331 München, Germany")
    doc.add_paragraph("")
    doc.add_paragraph("Dear Herr Müller and Frau Schäfer,")
    doc.add_paragraph("")
    doc.add_paragraph(
        "This letter confirms our engagement to represent Müller & Schäfer GmbH "
        "in connection with the acquisition of Brückner Technologien AG. "
        "All communications shall be in English unless otherwise agreed."
    )
    doc.add_paragraph("")
    doc.add_paragraph("DRAFT 2 — subject to partner review")
    doc.save(str(path))
    print(f"CREATED: {path.name}")


# ---------------------------------------------------------------------------
# 12. Scanned PDFs (VG-2 OCR pipeline fixtures) — image-only pages, no text
# layer, so extractPdfText reports scanned: true and the OCR pipeline runs.
# ---------------------------------------------------------------------------

# The unique quotable sentence planted in scanned-filing-stamped.pdf page 2.
# Tests cite this constant verbatim (tests/unit/ocr-pipeline.test.ts and the
# Task 14 native retrieval run) — change it ONLY together with those tests.
SCANNED_FILING_SENTENCE = (
    "Defendant's motion to compel production of the September audit file is DENIED."
)


def _scan_font(names, size):
    """Best-effort truetype font for the scanned pages. Tries the DejaVu family
    (preinstalled on the CI/dev Linux rigs); falls back to PIL's default font so
    the generator still runs anywhere, just with a less realistic page."""
    candidates = []
    for name in names:
        candidates.append(f"/usr/share/fonts/truetype/dejavu/{name}.ttf")
        candidates.append(name + ".ttf")  # cwd / font-path lookup
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except OSError:
            continue
    try:
        return ImageFont.load_default(size=size)  # Pillow >= 10.1
    except TypeError:
        return ImageFont.load_default()


def _draw_centered(draw, width, y, text, fnt):
    bbox = draw.textbbox((0, 0), text, font=fnt)
    draw.text(((width - bbox[2]) / 2, y), text, font=fnt, fill=0)


def make_scanned_filing_stamped():
    """scanned-filing-stamped.pdf — 2 image-only pages at 200 dpi: a clean
    motion-style court order with a FILED stamp box. Page 2 carries the unique
    quotable sentence (SCANNED_FILING_SENTENCE). Calibrated 2026-06-11 against
    tesseract 5.3.4: mean word confidence ~95 (well above the OCR_LOW_CONFIDENCE
    = 60 disclosure threshold — this is the GOOD scan)."""
    if not HAS_PIL:
        print("SKIP: scanned-filing-stamped.pdf (no Pillow)")
        return
    path = out("scanned-filing-stamped.pdf")
    w, h = 1700, 2200  # 8.5 x 11 in at 200 dpi
    serif = _scan_font(["DejaVuSerif"], 34)
    serif_bold = _scan_font(["DejaVuSerif-Bold"], 36)
    stamp_font = _scan_font(["DejaVuSans-Bold"], 30)
    leading = 56

    def lines_block(draw, x, y, lines):
        for line in lines:
            draw.text((x, y), line, font=serif, fill=0)
            y += leading
        return y

    # Page 1 — caption, stamp, the motion background.
    p1 = Image.new("L", (w, h), 255)
    d1 = ImageDraw.Draw(p1)
    y = 150
    _draw_centered(d1, w, y, "UNITED STATES DISTRICT COURT", serif_bold)
    y += 54
    _draw_centered(d1, w, y, "SOUTHERN DISTRICT OF NEW YORK", serif_bold)
    y += 120
    y = lines_block(d1, 160, y, [
        "MARCUS JOHNSON,",
        "                         Plaintiff,",
        "",
        "            v.                                          Case No. 26-CV-0412",
        "",
        "NEXUS DYNAMICS CORP.,",
        "                         Defendant.",
    ])
    y += 80
    _draw_centered(d1, w, y, "ORDER ON DEFENDANT'S MOTION TO COMPEL", serif_bold)
    y += 110
    lines_block(d1, 160, y, [
        "Before the Court is Defendant's motion to compel production of",
        "the September audit file and the related custodian communications.",
        "The Court has reviewed the parties' letter briefs, the deposition",
        "transcript of Marcus Johnson, and the record of the October 2025",
        "compliance review.",
        "",
        "Defendant argues the audit file is responsive to Request No. 12",
        "and that Plaintiff's objections were waived as untimely. Plaintiff",
        "responds that the file contains privileged work product prepared",
        "at the direction of counsel after the preservation notice issued.",
    ])
    # FILED stamp box overlay (top right). The em dash is part of the
    # stamped image artwork specified by the plan, not user-facing app copy.
    sx, sy = w - 560, 160
    d1.rectangle([sx, sy, sx + 410, sy + 150], outline=0, width=4)
    d1.text((sx + 26, sy + 16), "FILED — CLERK", font=stamp_font, fill=0)
    d1.text((sx + 26, sy + 60), "OF COURT", font=stamp_font, fill=0)
    d1.text((sx + 26, sy + 104), "JUN 02 2026", font=stamp_font, fill=0)

    # Page 2 — the ruling, carrying the planted quotable sentence.
    p2 = Image.new("L", (w, h), 255)
    d2 = ImageDraw.Draw(p2)
    y = 170
    y = lines_block(d2, 160, y, [
        "Having weighed the parties' submissions under Rule 26(b)(1), the",
        "Court finds the September audit file was prepared in anticipation",
        "of litigation following the document preservation notice, and that",
        "Defendant has not shown substantial need sufficient to overcome",
        "the work product protection.",
        "",
        "Accordingly, IT IS HEREBY ORDERED that:",
        "",
    ])
    # Wrap the planted sentence honestly across two lines at this type size.
    y = lines_block(d2, 160, y, [
        "Defendant's motion to compel production of the September",
        "audit file is DENIED.",
        "",
        "The parties shall meet and confer on the remaining custodian",
        "list disputes and submit a joint status letter within fourteen",
        "days of this Order.",
        "",
        "SO ORDERED.",
        "",
        "Dated: June 2, 2026",
        "New York, New York",
        "",
        "                                        ____________________________",
        "                                        United States District Judge",
    ])
    p1.save(str(path), format="PDF", save_all=True, append_images=[p2], resolution=200.0)
    print(f"CREATED: {path.name} (2 image-only pages, 200 dpi)")


def make_scanned_fax_noisy():
    """scanned-fax-noisy.pdf — 1 image-only page engineered to land BELOW the
    OCR_LOW_CONFIDENCE = 60 disclosure threshold. The Task 6 spike measured the
    plan's first recipe (2% noise + 2.5 deg) at confidence 63-66 — too high — so
    this uses heavier degradation: 5% salt-and-pepper noise, 3 deg rotation, and
    a 0.8 px gaussian blur. Calibrated 2026-06-11 against tesseract 5.3.4:
    mean word confidence ~53 (and tesseract-wasm measured ~2.5 LOWER than the
    native CLI on noisy pages in the spike). Deterministic via seed 42."""
    if not HAS_PIL:
        print("SKIP: scanned-fax-noisy.pdf (no Pillow)")
        return
    import random as _random
    rng = _random.Random(42)
    path = out("scanned-fax-noisy.pdf")
    w, h = 1275, 1650  # 8.5 x 11 in at 150 dpi
    mono = _scan_font(["DejaVuSansMono"], 22)
    small = _scan_font(["DejaVuSansMono"], 18)
    img = Image.new("L", (w, h), 255)
    d = ImageDraw.Draw(img)
    d.text(
        (60, 40),
        "FROM: MARCHETTI & ASSOCIATES   FAX 212-555-0142   06/02/2026 14:31   P.01/01",
        font=small,
        fill=0,
    )
    d.line([(60, 75), (w - 60, 75)], fill=0, width=2)
    body = [
        "RE: Johnson v. Nexus Dynamics Corp., Case No. 26-CV-0412",
        "",
        "Counsel:",
        "",
        "Per our call this morning, the September audit file remains",
        "outstanding. Your client's production is now six weeks late.",
        "We will move to compel on Thursday unless the file and the",
        "related custodian list are produced by close of business.",
        "",
        "The court reporter confirmed the deposition transcript for",
        "Marcus Johnson is final. Exhibit 14 was marked and admitted.",
        "",
        "Govern yourselves accordingly.",
        "",
        "Diane Marchetti, Esq.",
    ]
    y = 130
    for line in body:
        d.text((90, y), line, font=mono, fill=0)
        y += 36
    # Degradation: slight skew, soft strokes, then salt-and-pepper.
    img = img.rotate(3.0, expand=False, fillcolor=255, resample=Image.BICUBIC)
    img = img.filter(ImageFilter.GaussianBlur(0.8))
    px = img.load()
    flipped = int(w * h * 0.05)
    for _ in range(flipped):
        x = rng.randrange(w)
        yy = rng.randrange(h)
        px[x, yy] = 0 if rng.random() < 0.5 else 255
    img.save(str(path), format="PDF", resolution=150.0)
    print(f"CREATED: {path.name} (1 image-only page, degraded below confidence 60)")


# ---------------------------------------------------------------------------
# 11. matter-b-acme/ — 3 distinct files (Acme Corp. contract dispute)
# ---------------------------------------------------------------------------
def make_matter_b():
    # matter-b-1: intake memo
    if HAS_DOCX:
        path = out_b("intake-memo-acme.docx")
        doc = DocxDocument()
        doc.add_heading("CLIENT INTAKE MEMORANDUM", 0)
        doc.add_paragraph("Client: Acme Corporation (Wile E. Coyote, CEO)")
        doc.add_paragraph("Opposing Party: Road Runner Logistics LLC")
        doc.add_paragraph("Matter: Breach of Contract — Supply Chain Disruption")
        doc.add_paragraph("Intake Date: June 10, 2026")
        doc.add_paragraph("")
        doc.add_heading("Background", 1)
        doc.add_paragraph(
            "Acme Corporation entered into a 24-month supply agreement with "
            "Road Runner Logistics LLC in January 2025. Road Runner failed to "
            "deliver 14 of 24 monthly shipments on schedule. Acme suffered "
            "$2.4 million in downstream losses due to production shutdowns."
        )
        doc.add_heading("Claims", 1)
        doc.add_paragraph("1. Breach of contract (failure to perform)")
        doc.add_paragraph("2. Breach of implied covenant of good faith and fair dealing")
        doc.add_paragraph("3. Consequential damages for production shutdown losses")
        doc.save(str(path))
        print(f"CREATED: matter-b-acme/{path.name}")

    # matter-b-2: contract
    path2 = out_b("acme-supply-agreement.txt")
    path2.write_text("""\
SUPPLY AGREEMENT
Acme Corporation and Road Runner Logistics LLC
Effective January 1, 2025

1. PARTIES. Acme Corporation ("Buyer") and Road Runner Logistics LLC ("Supplier").

2. SERVICES. Supplier shall deliver 500 units of Widget Model X per month,
   on or before the 5th business day of each month.

3. PRICE. Buyer shall pay $48,000 per monthly shipment, net 30 days.

4. LIQUIDATED DAMAGES. For each day of delay beyond the delivery deadline,
   Supplier shall pay $2,400 as liquidated damages (not a penalty).

5. TERM. Twenty-four (24) months from the Effective Date.

6. GOVERNING LAW. State of Nevada.

SIGNED:
Acme Corporation: _______________
Road Runner Logistics LLC: _______________
""", encoding="utf-8")
    print(f"CREATED: matter-b-acme/{path2.name}")

    # matter-b-3: damages summary
    if HAS_XLSX:
        path3 = out_b("acme-damages-summary.xlsx")
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Acme Damages"
        ws.append(["Month", "Shipment Due", "Delivered", "Delay (days)", "Liquidated Damages"])
        delays = [0, 0, 3, 0, 12, 7, 0, 0, 22, 0, 0, 18, 5, 0, 30, 9, 0, 0, 0, 14, 0, 0, 8, 0]
        for i, d in enumerate(delays):
            month = f"2025-{(i % 12) + 1:02d}"
            delivered = "YES" if d == 0 else f"LATE +{d}d"
            ld = d * 2400 if d > 0 else 0
            ws.append([month, "5th business day", delivered, d, ld])
        ws.append(["", "", "", "TOTAL", f"=SUM(E2:E{len(delays)+1})"])
        wb.save(str(path3))
        print(f"CREATED: matter-b-acme/{path3.name}")
    else:
        path3b = out_b("acme-damages-summary.md")
        path3b.write_text("# Acme Damages Summary\n\nSee legal team for full spreadsheet.\n", encoding="utf-8")
        print(f"CREATED: matter-b-acme/{path3b.name}")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    # Each maker has a stable name so a subset can be regenerated without
    # rewriting the other committed binaries (docx/xlsx/pptx zips embed
    # timestamps — a full rerun dirties every fixture in git).
    MAKERS = [
        ("engagement-letter", make_engagement_letter_tracked),
        ("deposition-transcript", make_deposition_transcript),
        ("weston-transcript", make_weston_certified_transcript),
        ("incident-summary", make_incident_summary),
        ("services-agreement", make_services_agreement),
        ("scanned-exhibit", make_scanned_exhibit),
        ("damages-model", make_damages_model),
        ("exhibit-deck", make_exhibit_deck),
        ("huge-notes", make_huge_notes),
        ("empty-docx", make_empty_docx),
        ("unicode-docx", make_unicode_docx),
        ("matter-b", make_matter_b),
        ("scanned-filing", make_scanned_filing_stamped),
        ("scanned-fax", make_scanned_fax_noisy),
    ]
    only = set(sys.argv[1:])
    unknown = only - {name for name, _ in MAKERS}
    if unknown:
        print(f"Unknown maker name(s): {sorted(unknown)}", file=sys.stderr)
        sys.exit(1)
    print(f"Generating fixtures in: {CORPUS_DIR}")
    for name, maker in MAKERS:
        if not only or name in only:
            maker()
    print("\nDone.")
